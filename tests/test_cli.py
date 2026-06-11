import json
import os
import subprocess
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches
from pptx_cli.session import (
    cleanup_stale_session_state,
    resolve_repo_root,
    resolve_state_file_path,
    save_session_state,
)
from pptx_cli.template_ops import resolve_template_root

ROOT = Path(__file__).resolve().parents[1]
CLI = ROOT / "pptxcli"


def run_cli(*args: str, env: dict[str, str] | None = None) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [str(CLI), *args],
        check=False,
        capture_output=True,
        text=True,
        env=env,
    )


class CliSmokeTest(unittest.TestCase):
    def test_help_runs(self) -> None:
        result = run_cli("--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("pptxcli", result.stdout)

    def test_show_help_runs(self) -> None:
        result = run_cli("show", "--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("--annotate", result.stdout)

    def test_inspect_help_runs(self) -> None:
        result = run_cli("inspect", "--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("--slide", result.stdout)

    def test_edit_help_runs(self) -> None:
        result = run_cli("edit", "--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("create", result.stdout)

    def test_edit_create_help_runs(self) -> None:
        result = run_cli("edit", "create", "--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("--output", result.stdout)

    def test_edit_show_template_help_runs(self) -> None:
        result = run_cli("edit", "show_template", "--help")
        self.assertEqual(result.returncode, 0)
        self.assertIn("--slide", result.stdout)

    def test_demo_form_outputs_json(self) -> None:
        result = run_cli("demo", "form")
        self.assertEqual(result.returncode, 0)
        payload = json.loads(result.stdout)
        self.assertEqual(payload["template_dir"], "./templates/demo_template")

    def test_session_init_inspect_finish_roundtrip(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            pptx_path = tmp_dir / "demo.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            text_box.text_frame.text = "Session Ready"
            presentation.save(pptx_path)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)

            init_result = run_cli("init", "--origin_file", str(pptx_path), env=env)
            self.assertEqual(init_result.returncode, 0, init_result.stderr)
            init_payload = json.loads(init_result.stdout)
            self.assertEqual(init_payload["status"], "ready")
            self.assertEqual(init_payload["origin_file"], str(pptx_path.resolve()))

            inspect_result = run_cli("inspect", "--slide", "0", env=env)
            self.assertEqual(inspect_result.returncode, 0, inspect_result.stderr)
            inspect_payload = json.loads(inspect_result.stdout)
            self.assertEqual(inspect_payload["command"], "inspect")
            self.assertEqual(len(inspect_payload["objects"]), 1)
            self.assertEqual(inspect_payload["objects"][0]["object_type"], "text")
            self.assertEqual(inspect_payload["input"], str(pptx_path.resolve()))

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)
            finish_payload = json.loads(finish_result.stdout)
            self.assertEqual(finish_payload["status"], "stopped")
            self.assertFalse(state_file.exists())

    def test_cleanup_stale_session_state_handles_request_timeout(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            save_session_state(
                state_file,
                {
                    "server_url": "http://127.0.0.1:65535",
                    "mode": "edit_ppt",
                },
            )

            with patch("pptx_cli.session.request.urlopen", side_effect=TimeoutError("timed out")):
                cleaned = cleanup_stale_session_state(state_file)

            self.assertTrue(cleaned)
            self.assertFalse(state_file.exists())

    def test_resolve_paths_follow_repo_root(self) -> None:
        with patch.dict(os.environ, {}, clear=True):
            self.assertEqual(resolve_repo_root(), ROOT)
            self.assertEqual(resolve_state_file_path(), ROOT / ".pptxcli-session.json")
            self.assertEqual(resolve_template_root(), ROOT / "templates")


if __name__ == "__main__":
    unittest.main()
