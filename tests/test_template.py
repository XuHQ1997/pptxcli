import json
import io
import os
import subprocess
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_cli.container_layout import (
    flatten_resolved_leaves,
    parse_content_spec,
    solve_content_layout,
)
from pptx_cli.inspect import inspect_slide_objects
from pptx_cli.session import save_session_state
from pptx_cli.template_ops import save_template_package

ROOT = Path(__file__).resolve().parents[1]
CLI = ROOT / "pptxcli"


def run_cli(*args: str, env: dict[str, str] | None = None) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [str(CLI), *args],
        check=False,
        capture_output=True,
        text=True,
        env=env,
    )


def collect_slide_texts(presentation: Presentation) -> list[str]:
    slide_texts: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = " ".join(shape.text.split())
                if text:
                    parts.append(text)
        slide_texts.append(" | ".join(parts))
    return slide_texts


def collect_slide_image_sizes(presentation: Presentation) -> list[list[tuple[int, int]]]:
    image_sizes: list[list[tuple[int, int]]] = []
    for slide in presentation.slides:
        current_slide_sizes: list[tuple[int, int]] = []
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                with Image.open(io.BytesIO(shape.image.blob)) as image:
                    current_slide_sizes.append(image.size)
        image_sizes.append(current_slide_sizes)
    return image_sizes


def collect_first_run_formats(presentation: Presentation) -> list[list[dict[str, object]]]:
    slide_formats: list[list[dict[str, object]]] = []
    for slide in presentation.slides:
        current_slide_formats: list[dict[str, object]] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if not paragraph.runs:
                    continue
                run = paragraph.runs[0]
                current_slide_formats.append(
                    {
                        "text": run.text,
                        "font_name": run.font.name,
                        "font_size": run.font.size.pt if run.font.size is not None else None,
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                    }
                )
        slide_formats.append(current_slide_formats)
    return slide_formats


def collect_paragraph_run_counts(presentation: Presentation) -> list[list[int]]:
    slide_counts: list[list[int]] = []
    for slide in presentation.slides:
        current_slide_counts: list[int] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                current_slide_counts.append(len(paragraph.runs))
        slide_counts.append(current_slide_counts)
    return slide_counts


def collect_slide_embed_relationships(pptx_path: Path) -> list[dict[str, set[str]]]:
    slide_relationships: list[dict[str, set[str]]] = []
    with zipfile.ZipFile(pptx_path) as archive:
        slide_indexes = sorted(
            int(path.stem.replace("slide", ""))
            for path in (Path(name) for name in archive.namelist())
            if path.parent == Path("ppt/slides") and path.suffix == ".xml" and path.stem.startswith("slide")
        )
        for slide_index in slide_indexes:
            slide_xml = archive.read(f"ppt/slides/slide{slide_index}.xml").decode(
                "utf-8",
                errors="ignore",
            )
            rels_xml = archive.read(
                f"ppt/slides/_rels/slide{slide_index}.xml.rels"
            ).decode("utf-8", errors="ignore")
            embed_ids = set()
            rel_ids = set()
            for marker in ('r:embed="',):
                start = 0
                while True:
                    position = slide_xml.find(marker, start)
                    if position < 0:
                        break
                    value_start = position + len(marker)
                    value_end = slide_xml.find('"', value_start)
                    embed_ids.add(slide_xml[value_start:value_end])
                    start = value_end + 1
            for marker in ('Id="',):
                start = 0
                while True:
                    position = rels_xml.find(marker, start)
                    if position < 0:
                        break
                    value_start = position + len(marker)
                    value_end = rels_xml.find('"', value_start)
                    rel_ids.add(rels_xml[value_start:value_end])
                    start = value_end + 1
            slide_relationships.append({"embed_ids": embed_ids, "rel_ids": rel_ids})
    return slide_relationships


def inspect_template_objects(pptx_path: Path, slide_index: int) -> list[object]:
    return list(inspect_slide_objects(pptx_path, slide_index).objects)


class TemplateWorkflowTest(unittest.TestCase):
    def test_container_layout_solves_nested_vertical_horizontal_tree(self) -> None:
        content_payload = {
            "layout": "vertical",
            "gap": 0.05,
            "children": [
                {
                    "type": "text",
                    "name": "title",
                    "ratio": 1,
                    "text": "Quarterly Review",
                },
                {
                    "layout": "horizontal",
                    "name": "body",
                    "ratio": 3,
                    "gap": 0.05,
                    "children": [
                        {
                            "type": "text",
                            "name": "left",
                            "ratio": 2,
                            "text": "Left Column",
                        },
                        {
                            "type": "text",
                            "name": "right",
                            "ratio": 1,
                            "text": "Right Column",
                        },
                    ],
                },
            ],
        }

        tree = parse_content_spec(
            json.dumps(content_payload),
            slide_width=1000,
            slide_height=800,
        )
        resolved = solve_content_layout(tree, slide_width=1000, slide_height=800)
        leaves = {leaf.name: leaf.bbox for leaf in flatten_resolved_leaves(resolved)}

        self.assertEqual(resolved.bbox.to_dict(), {"x": 0, "y": 0, "w": 1000, "h": 800})
        self.assertEqual(leaves["title"].to_dict(), {"x": 0, "y": 0, "w": 1000, "h": 190})
        self.assertEqual(leaves["left"].to_dict(), {"x": 0, "y": 230, "w": 633, "h": 570})
        self.assertEqual(leaves["right"].to_dict(), {"x": 683, "y": 230, "w": 317, "h": 570})

    def test_container_layout_ratio_array_uses_axis_specific_values(self) -> None:
        tree = parse_content_spec(
            json.dumps(
                {
                    "layout": "vertical",
                    "children": [
                        {"type": "text", "name": "top", "ratio": [1, 1], "text": "Top"},
                        {
                            "layout": "horizontal",
                            "name": "bottom",
                            "ratio": [1, 3],
                            "children": [
                                {"type": "text", "name": "left", "ratio": [3, 1], "text": "Left"},
                                {"type": "text", "name": "right", "ratio": [1, 1], "text": "Right"},
                            ],
                        },
                    ],
                }
            ),
            slide_width=1000,
            slide_height=800,
        )
        resolved = solve_content_layout(tree, slide_width=1000, slide_height=800)
        leaves = {leaf.name: leaf.bbox for leaf in flatten_resolved_leaves(resolved)}

        self.assertEqual(leaves["top"].to_dict(), {"x": 0, "y": 0, "w": 1000, "h": 200})
        self.assertEqual(leaves["left"].to_dict(), {"x": 0, "y": 200, "w": 750, "h": 600})
        self.assertEqual(leaves["right"].to_dict(), {"x": 750, "y": 200, "w": 250, "h": 600})

    def test_container_layout_grid_uses_max_axis_ratio_per_track(self) -> None:
        tree = parse_content_spec(
            json.dumps(
                {
                    "layout": "grid",
                    "columns": 2,
                    "children": [
                        {"type": "text", "name": "a1", "ratio": [2, 1], "text": "A1"},
                        {"type": "text", "name": "a2", "ratio": [1, 1], "text": "A2"},
                        {"type": "text", "name": "b1", "ratio": [3, 2], "text": "B1"},
                        {"type": "text", "name": "b2", "ratio": [4, 5], "text": "B2"},
                    ],
                }
            ),
            slide_width=1000,
            slide_height=800,
        )
        resolved = solve_content_layout(tree, slide_width=1000, slide_height=800)
        leaves = {leaf.name: leaf.bbox for leaf in flatten_resolved_leaves(resolved)}

        self.assertEqual(leaves["a1"].to_dict(), {"x": 0, "y": 0, "w": 429, "h": 133})
        self.assertEqual(leaves["a2"].to_dict(), {"x": 429, "y": 0, "w": 571, "h": 133})
        self.assertEqual(leaves["b1"].to_dict(), {"x": 0, "y": 133, "w": 429, "h": 667})
        self.assertEqual(leaves["b2"].to_dict(), {"x": 429, "y": 133, "w": 571, "h": 667})

    def test_template_create_add_slide_save_roundtrip(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            image_path = tmp_dir / "hero.png"
            Image.new("RGB", (80, 80), "green").save(image_path)

            presentation = Presentation()
            slide_0 = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_0.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )

            slide_1 = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_1.shapes.add_picture(
                str(image_path),
                Inches(1),
                Inches(1.5),
                Inches(2),
                Inches(2),
            )

            slide_2 = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Summary Page"
            )
            presentation.save(pptx_path)

            slide_0_objects = inspect_template_objects(pptx_path, 0)
            slide_2_objects = inspect_template_objects(pptx_path, 2)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            create_result = run_cli(
                "template",
                "create",
                "--from",
                str(pptx_path),
                "--name",
                "demo_template.pptx",
                env=env,
            )
            self.assertEqual(create_result.returncode, 0, create_result.stderr)
            create_payload = json.loads(create_result.stdout)
            self.assertEqual(create_payload["template_name"], "demo_template")
            draft_path = Path(create_payload["draft_path"])
            self.assertTrue(draft_path.exists())

            add_slide_2_result = run_cli(
                "template",
                "add_slide",
                "--slide",
                "2",
                "-f",
                f"{slide_2_objects[0].index}:summary title",
                env=env,
            )
            self.assertEqual(add_slide_2_result.returncode, 0, add_slide_2_result.stderr)

            add_slide_0_result = run_cli(
                "template",
                "add_slide",
                "--slide",
                "0",
                "-f",
                f"{slide_0_objects[0].index}:hero title",
                env=env,
            )
            self.assertEqual(add_slide_0_result.returncode, 0, add_slide_0_result.stderr)

            save_result = run_cli("template", "save", env=env)
            self.assertEqual(save_result.returncode, 0, save_result.stderr)
            save_payload = json.loads(save_result.stdout)
            manifest_path = Path(save_payload["manifest_path"])
            template_pptx_path = Path(save_payload["template_pptx_path"])

            self.assertTrue(manifest_path.exists())
            self.assertTrue(template_pptx_path.exists())

            draft_payload = json.loads(draft_path.read_text(encoding="utf-8"))
            self.assertEqual(
                [slide["source_slide_index"] for slide in draft_payload["slides"]],
                [2, 0],
            )

            manifest_payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            self.assertEqual(manifest_payload["template_name"], "demo_template")
            self.assertEqual(manifest_payload["slide_count"], 2)
            self.assertEqual(manifest_payload["field_count"], 2)
            self.assertEqual(
                [slide["source_slide_index"] for slide in manifest_payload["slides"]],
                [2, 0],
            )
            self.assertEqual(
                [slide["slide_name"] for slide in manifest_payload["slides"]],
                ["slide_2", "slide_0"],
            )
            self.assertEqual(
                [field["description"] for field in manifest_payload["fields"]],
                ["summary title", "hero title"],
            )
            self.assertEqual(
                [field["index"] for field in manifest_payload["fields"]],
                [1, 1],
            )

            template_presentation = Presentation(str(template_pptx_path))
            self.assertEqual(len(template_presentation.slides), 2)
            self.assertEqual(
                collect_slide_texts(template_presentation),
                ["Summary Page", "Hero Title"],
            )

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_template_save_cleans_repo_preview_dir(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            repo_root = tmp_dir / "repo"
            preview_dir = repo_root / "preview"
            template_root = tmp_dir / "template-store"
            state_file = tmp_dir / ".pptxcli-session.json"
            pptx_path = tmp_dir / "demo.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            presentation.save(pptx_path)

            template_root.mkdir(parents=True, exist_ok=True)
            draft_path = template_root / "demo_preview_cleanup.json"
            draft_path.write_text(
                json.dumps(
                    {
                        "draft_version": 1,
                        "template_name": "demo_preview_cleanup",
                        "origin_file": str(pptx_path.resolve()),
                        "created_at": "2026-01-01T00:00:00+00:00",
                        "slides": [
                            {
                                "slide_name": "slide_0",
                                "source_slide_index": 0,
                                "slide_size": None,
                                "fields": [],
                                "content_shapes": [],
                                "content_area": None,
                                "added_at": "2026-01-01T00:00:00+00:00",
                            }
                        ],
                    },
                    indent=2,
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            save_session_state(
                state_file,
                {
                    "mode": "template_extract",
                    "active_template": "demo_preview_cleanup",
                },
            )

            preview_dir.mkdir(parents=True, exist_ok=True)
            (preview_dir / "demo.slide-0.annotated.png").write_bytes(b"fake-image")

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            with patch.dict(os.environ, env, clear=False):
                with patch("pptx_cli.template_ops.resolve_repo_root", return_value=repo_root):
                    payload = save_template_package()

            self.assertEqual(payload["status"], "saved")
            self.assertFalse(preview_dir.exists())

    def test_edit_fill_template_appends_multiple_slides_and_save(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            source_image_path = tmp_dir / "source.png"
            replacement_image_path = tmp_dir / "replacement.png"
            output_path = tmp_dir / "final-output.pptx"

            Image.new("RGB", (80, 80), "green").save(source_image_path)
            Image.new("RGB", (120, 60), "red").save(replacement_image_path)

            presentation = Presentation()
            slide_0 = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_0.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )

            slide_1 = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide_1.shapes.add_picture(
                str(source_image_path),
                Inches(1),
                Inches(1.5),
                Inches(2),
                Inches(2),
            )
            presentation.save(pptx_path)

            slide_0_objects = inspect_template_objects(pptx_path, 0)
            slide_1_objects = inspect_template_objects(pptx_path, 1)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_fill",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_0_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "1",
                    "-f",
                    f"{slide_1_objects[0].index}:cover image",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)

            create_edit_result = run_cli(
                "edit",
                "create",
                "--output",
                str(output_path),
                "--template",
                "demo_fill",
                env=env,
            )
            self.assertEqual(create_edit_result.returncode, 0, create_edit_result.stderr)
            create_edit_payload = json.loads(create_edit_result.stdout)
            self.assertEqual(create_edit_payload["status"], "created")
            self.assertEqual(create_edit_payload["slide_count"], 0)

            state_payload = json.loads(state_file.read_text(encoding="utf-8"))
            self.assertEqual(state_payload["mode"], "edit_ppt")
            self.assertEqual(state_payload["active_template"], "demo_fill")
            self.assertEqual(
                Path(state_payload["edit_context"]["output_path"]),
                output_path.resolve(),
            )

            fill_text_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan",
                env=env,
            )
            self.assertEqual(fill_text_result.returncode, 0, fill_text_result.stderr)
            fill_text_payload = json.loads(fill_text_result.stdout)
            self.assertEqual(fill_text_payload["slide_count"], 1)
            self.assertEqual(fill_text_payload["slide"], "slide_0")

            fill_image_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "1",
                "-f",
                f"1:{replacement_image_path}",
                env=env,
            )
            self.assertEqual(fill_image_result.returncode, 0, fill_image_result.stderr)
            fill_image_payload = json.loads(fill_image_result.stdout)
            self.assertEqual(fill_image_payload["slide_count"], 2)
            self.assertEqual(fill_image_payload["slide"], "slide_1")

            save_result = run_cli("edit", "save", env=env)
            self.assertEqual(save_result.returncode, 0, save_result.stderr)
            save_payload = json.loads(save_result.stdout)
            self.assertEqual(save_payload["status"], "saved")
            self.assertEqual(Path(save_payload["output_path"]), output_path.resolve())

            state_payload = json.loads(state_file.read_text(encoding="utf-8"))
            self.assertEqual(state_payload["mode"], "template_extract")
            self.assertIsNone(state_payload.get("edit_context"))

            output_presentation = Presentation(str(output_path))
            self.assertEqual(len(output_presentation.slides), 2)
            self.assertEqual(
                collect_slide_texts(output_presentation),
                ["Launch Plan", ""],
            )
            self.assertEqual(
                collect_slide_image_sizes(output_presentation),
                [[], [(120, 60)]],
            )

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_show_template_lists_manifest_fields(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "filled-output.pptx"
            image_path = tmp_dir / "hero.png"

            Image.new("RGB", (80, 80), "green").save(image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            slide.shapes.add_picture(
                str(image_path),
                Inches(1),
                Inches(2),
                Inches(2),
                Inches(2),
            )
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)
            text_candidate = next(item for item in slide_objects if item.object_type == "text")
            image_candidate = next(item for item in slide_objects if item.object_type == "image")

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_show",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{text_candidate.index}:hero title",
                    "-f",
                    f"{image_candidate.index}:hero image",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "demo_show",
                    env=env,
                ).returncode,
                0,
            )

            show_result = run_cli(
                "edit",
                "show_template",
                "--slide",
                "0",
                env=env,
            )
            self.assertEqual(show_result.returncode, 0, show_result.stderr)
            show_payload = json.loads(show_result.stdout)
            self.assertEqual(show_payload["command"], "edit show_template")
            self.assertEqual(show_payload["mode"], "edit_ppt")
            self.assertEqual(show_payload["template_name"], "demo_show")
            self.assertEqual(show_payload["slide"]["slide_name"], "slide_0")
            self.assertEqual(show_payload["slide"]["field_count"], 2)
            self.assertEqual(
                [field["index"] for field in show_payload["slide"]["fields"]],
                [1, 2],
            )
            self.assertEqual(
                [field["description"] for field in show_payload["slide"]["fields"]],
                ["hero title", "hero image"],
            )
            self.assertEqual(
                [field["type"] for field in show_payload["slide"]["fields"]],
                ["text", "image"],
            )

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_preserves_text_formatting(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "styled-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            paragraph = text_box.text_frame.paragraphs[0]
            paragraph.clear()
            run = paragraph.add_run()
            run.text = "Hero"
            run.font.name = "Arial"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.italic = True
            accent_run = paragraph.add_run()
            accent_run.text = " Title"
            accent_run.font.name = "Courier New"
            accent_run.font.size = Pt(16)
            accent_run.font.bold = False
            accent_run.font.italic = False
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_style",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "demo_style",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan",
                env=env,
            )
            self.assertEqual(fill_result.returncode, 0, fill_result.stderr)
            self.assertEqual(run_cli("edit", "save", env=env).returncode, 0)

            output_presentation = Presentation(str(output_path))
            self.assertEqual(collect_slide_texts(output_presentation), ["Launch Plan"])
            run_formats = collect_first_run_formats(output_presentation)
            run_counts = collect_paragraph_run_counts(output_presentation)
            self.assertEqual(run_formats[0][0]["text"], "Launch Plan")
            self.assertEqual(run_formats[0][0]["font_name"], "Arial")
            self.assertEqual(run_formats[0][0]["font_size"], 28.0)
            self.assertTrue(run_formats[0][0]["bold"])
            self.assertTrue(run_formats[0][0]["italic"])
            self.assertEqual(run_counts[0], [1])

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_preserves_multiline_text_formatting(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "styled-multiline-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
            first_paragraph = text_box.text_frame.paragraphs[0]
            first_paragraph.clear()
            first_run = first_paragraph.add_run()
            first_run.text = "Hero"
            first_run.font.name = "Arial"
            first_run.font.size = Pt(28)
            first_run.font.bold = True
            first_paragraph.add_run().text = " Title"

            second_paragraph = text_box.text_frame.add_paragraph()
            second_run = second_paragraph.add_run()
            second_run.text = "Summary"
            second_run.font.name = "Courier New"
            second_run.font.size = Pt(16)
            second_run.font.italic = True
            second_paragraph.add_run().text = " Extra"
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_multiline_style",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero body",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "demo_multiline_style",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan\nKey Summary",
                env=env,
            )
            self.assertEqual(fill_result.returncode, 0, fill_result.stderr)
            self.assertEqual(run_cli("edit", "save", env=env).returncode, 0)

            output_presentation = Presentation(str(output_path))
            self.assertEqual(collect_slide_texts(output_presentation), ["Launch Plan Key Summary"])
            run_formats = collect_first_run_formats(output_presentation)
            run_counts = collect_paragraph_run_counts(output_presentation)
            self.assertEqual(run_formats[0][0]["text"], "Launch Plan")
            self.assertEqual(run_formats[0][0]["font_name"], "Arial")
            self.assertEqual(run_formats[0][0]["font_size"], 28.0)
            self.assertTrue(run_formats[0][0]["bold"])
            self.assertEqual(run_formats[0][1]["text"], "Key Summary")
            self.assertEqual(run_formats[0][1]["font_name"], "Courier New")
            self.assertEqual(run_formats[0][1]["font_size"], 16.0)
            self.assertTrue(run_formats[0][1]["italic"])
            self.assertEqual(run_counts[0], [1, 1])

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_preserves_group_images_and_slide_layout(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "group-layout-demo.pptx"
            group_image_path = tmp_dir / "group.png"
            replacement_image_path = tmp_dir / "replacement.png"
            output_path = tmp_dir / "group-layout-output.pptx"

            Image.new("RGB", (80, 80), "blue").save(group_image_path)
            Image.new("RGB", (120, 120), "red").save(replacement_image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            title_shape = slide.shapes.title
            if title_shape is None:
                self.fail("expected title placeholder on slide layout 1")
            title_shape.text = "Section Title"
            content_shape = slide.placeholders[1]
            content_shape.text = "Body Copy"
            slide.shapes.add_picture(
                str(replacement_image_path),
                Inches(5.5),
                Inches(1.5),
                Inches(1.5),
                Inches(1.5),
            )
            group_shape = slide.shapes.add_group_shape()
            group_shape.shapes.add_picture(
                str(group_image_path),
                Inches(0.2),
                Inches(0.2),
                Inches(0.5),
                Inches(0.5),
            )
            group_shape.shapes.add_textbox(
                Inches(0.8),
                Inches(0.2),
                Inches(1.2),
                Inches(0.5),
            ).text = "Badge"
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)
            text_candidate = next(item for item in slide_objects if item.object_type == "text")
            image_candidate = next(item for item in slide_objects if item.object_type == "image")
            source_layout_partname = str(
                Presentation(str(pptx_path)).slides[0].slide_layout.part.partname
            )

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "group_layout",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{text_candidate.index}:title",
                    "-f",
                    f"{image_candidate.index}:hero",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "group_layout",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Updated Title",
                "-f",
                f"2:{replacement_image_path}",
                env=env,
            )
            self.assertEqual(fill_result.returncode, 0, fill_result.stderr)
            self.assertEqual(run_cli("edit", "save", env=env).returncode, 0)

            output_presentation = Presentation(str(output_path))
            output_slide = output_presentation.slides[0]
            self.assertEqual(str(output_slide.slide_layout.part.partname), source_layout_partname)

            group_shapes = [
                shape
                for shape in output_slide.shapes
                if getattr(getattr(shape, "shape_type", None), "name", None) == "GROUP"
            ]
            self.assertEqual(len(group_shapes), 1)
            self.assertTrue(any(hasattr(shape, "image") for shape in group_shapes[0].shapes))

            slide_relationships = collect_slide_embed_relationships(output_path)
            self.assertTrue(slide_relationships[0]["embed_ids"])
            self.assertTrue(slide_relationships[0]["embed_ids"].issubset(slide_relationships[0]["rel_ids"]))

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_supports_nested_container_content(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            content_image_path = tmp_dir / "content.png"
            output_path = tmp_dir / "content-output.pptx"

            Image.new("RGB", (120, 60), "purple").save(content_image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5), Inches(0.8)).text_frame.text = "Template Title"
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(3.0))
            content_box.text_frame.text = "Sample Body"
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)
            title_object = next(item for item in slide_objects if getattr(item, "text", None) == "Template Title")
            content_object = next(item for item in slide_objects if getattr(item, "text", None) == "Sample Body")
            expected_content_bbox = {
                "x": content_object.bbox.x,
                "y": content_object.bbox.y,
                "w": content_object.bbox.w,
                "h": content_object.bbox.h,
            }

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "container_content",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{title_object.index}:page title",
                    "--content",
                    str(content_object.index),
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "container_content",
                    env=env,
                ).returncode,
                0,
            )

            content_payload = {
                "layout": "vertical",
                "gap": 0.05,
                "children": [
                    {
                        "type": "text",
                        "name": "agenda",
                        "ratio": 1,
                        "text": "Agenda",
                        "style": {"font_size": 20, "bold": True},
                    },
                    {
                        "layout": "horizontal",
                        "name": "body",
                        "ratio": 2,
                        "gap": 0.05,
                        "children": [
                            {
                                "type": "text",
                                "name": "summary",
                                "ratio": 2,
                                "text": "Summary Block",
                            },
                            {
                                "type": "image",
                                "name": "hero_image",
                                "ratio": 1,
                                "path": str(content_image_path),
                                "fit": "contain",
                            },
                        ],
                    },
                ],
            }

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan",
                "--content",
                json.dumps(content_payload),
                env=env,
            )
            self.assertEqual(fill_result.returncode, 0, fill_result.stderr)
            fill_payload = json.loads(fill_result.stdout)
            self.assertEqual(fill_payload["content_count"], 3)
            self.assertEqual(fill_payload["field_count"], 1)
            self.assertIsInstance(fill_payload["content_layout"], dict)
            self.assertEqual(len(fill_payload["rendered_content"]), 3)
            self.assertEqual(fill_payload["content_layout"]["bbox"], expected_content_bbox)

            self.assertEqual(run_cli("edit", "save", env=env).returncode, 0)

            manifest_path = template_root / "container_content" / "manifest.json"
            template_pptx_path = template_root / "container_content" / "template.pptx"
            manifest_payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            self.assertEqual(
                manifest_payload["slides"][0]["content_area"],
                expected_content_bbox,
            )
            template_presentation = Presentation(str(template_pptx_path))
            self.assertNotIn("Sample Body", collect_slide_texts(template_presentation)[0])

            output_presentation = Presentation(str(output_path))
            self.assertEqual(len(output_presentation.slides), 1)
            self.assertIn("Launch Plan", collect_slide_texts(output_presentation)[0])
            self.assertIn("Agenda", collect_slide_texts(output_presentation)[0])
            self.assertIn("Summary Block", collect_slide_texts(output_presentation)[0])
            self.assertEqual(collect_slide_image_sizes(output_presentation), [[(120, 60)]])

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_container_layout_gap_uses_ratio_value(self) -> None:
        tree = parse_content_spec(
            json.dumps(
                {
                    "layout": "vertical",
                    "gap": 0.05,
                    "children": [
                        {"type": "text", "text": "Top"},
                        {"type": "text", "text": "Bottom"},
                    ],
                }
            ),
            slide_width=1000,
            slide_height=800,
        )
        resolved = solve_content_layout(tree, slide_width=1000, slide_height=800)
        leaves = flatten_resolved_leaves(resolved)

        self.assertEqual(leaves[0].bbox.to_dict(), {"x": 0, "y": 0, "w": 1000, "h": 380})
        self.assertEqual(leaves[1].bbox.to_dict(), {"x": 0, "y": 420, "w": 1000, "h": 380})

    def test_container_layout_rejects_legacy_bbox_and_size_fields(self) -> None:
        with self.assertRaisesRegex(ValueError, "unsupported legacy container fields"):
            parse_content_spec(
                json.dumps(
                    {
                        "layout": "vertical",
                        "bbox": {"x": 0, "y": 0, "w": 100, "h": 100},
                        "children": [{"type": "text", "text": "Legacy"}],
                    }
                ),
                slide_width=1000,
                slide_height=800,
            )

    def test_edit_fill_template_requires_active_edit(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_text_fill",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan",
                env=env,
            )
            self.assertEqual(fill_result.returncode, 1)
            self.assertIn("no active edit draft", fill_result.stderr)

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_rejects_content_for_slide_without_content_area(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "title-only-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "title_only",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "title_only",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "1:Launch Plan",
                "--content",
                json.dumps(
                    {
                        "layout": "vertical",
                        "children": [
                            {
                                "type": "text",
                                "text": "Agenda",
                            }
                        ],
                    }
                ),
                env=env,
            )
            self.assertEqual(fill_result.returncode, 1)
            self.assertIn("does not define a content area", fill_result.stderr)

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_reports_missing_field(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "filled-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_error",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "demo_error",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                env=env,
            )
            self.assertNotEqual(fill_result.returncode, 0)
            self.assertIn("required", fill_result.stderr)

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)

    def test_edit_fill_template_reports_missing_field_index(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            state_file = tmp_dir / ".pptxcli-session.json"
            template_root = tmp_dir / "template-store"
            pptx_path = tmp_dir / "demo.pptx"
            output_path = tmp_dir / "filled-output.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.text = (
                "Hero Title"
            )
            presentation.save(pptx_path)

            slide_objects = inspect_template_objects(pptx_path, 0)

            env = os.environ.copy()
            env["PPTXCLI_STATE_FILE"] = str(state_file)
            env["PPTXCLI_TEMPLATE_ROOT"] = str(template_root)

            self.assertEqual(
                run_cli(
                    "template",
                    "create",
                    "--from",
                    str(pptx_path),
                    "--name",
                    "demo_error2",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(
                run_cli(
                    "template",
                    "add_slide",
                    "--slide",
                    "0",
                    "-f",
                    f"{slide_objects[0].index}:hero title",
                    env=env,
                ).returncode,
                0,
            )
            self.assertEqual(run_cli("template", "save", env=env).returncode, 0)
            self.assertEqual(
                run_cli(
                    "edit",
                    "create",
                    "--output",
                    str(output_path),
                    "--template",
                    "demo_error2",
                    env=env,
                ).returncode,
                0,
            )

            fill_result = run_cli(
                "edit",
                "fill_template",
                "--slide",
                "0",
                "-f",
                "99:ignored",
                env=env,
            )
            self.assertEqual(fill_result.returncode, 1)
            self.assertIn("does not exist", fill_result.stderr)

            finish_result = run_cli("finish", env=env)
            self.assertEqual(finish_result.returncode, 0, finish_result.stderr)


if __name__ == "__main__":
    unittest.main()
