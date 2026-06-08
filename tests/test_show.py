from __future__ import annotations

import json
import tempfile
import unittest
from contextlib import redirect_stdout
from io import StringIO
from pathlib import Path
from unittest.mock import patch

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches

from pptx_cli.cli import main
from pptx_cli.inspect import inspect_slide, inspect_slide_objects
from pptx_cli.models import BBox, Candidate
from pptx_cli.show import _fit_label_layout, annotate_candidates


class ShowCommandTest(unittest.TestCase):
    def test_inspect_slide_detects_text_and_image_candidates(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            pptx_path = tmp_dir / "demo.pptx"
            image_path = tmp_dir / "sample.png"
            Image.new("RGB", (80, 80), "blue").save(image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            text_box.text_frame.text = "Quarterly Review"
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(2), Inches(2))
            presentation.save(pptx_path)

            result = inspect_slide(pptx_path, 0)

            self.assertEqual(len(result.candidates), 2)
            self.assertEqual(result.candidates[0].kind, "text")
            self.assertEqual(result.candidates[0].role_hint, "title")
            self.assertEqual(result.candidates[1].kind, "image")

    def test_inspect_slide_objects_returns_all_shapes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            pptx_path = tmp_dir / "demo.pptx"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = (
                "Quarterly Review"
            )
            slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(2), Inches(1))
            presentation.save(pptx_path)

            result = inspect_slide_objects(pptx_path, 0)

            self.assertEqual(len(result.objects), 2)
            self.assertEqual(result.objects[0].object_type, "text")
            self.assertEqual(result.objects[1].object_type, "shape")

    def test_annotate_candidates_draws_boxes(self) -> None:
        base_image = Image.new("RGB", (400, 300), "white")
        candidates = [
            Candidate(
                index=1,
                candidate_id="slide-shape-1",
                shape_id=1,
                kind="text",
                role_hint="title",
                bbox=BBox(x=100, y=80, w=120, h=90),
                text_excerpt="Quarterly Review",
            )
        ]

        annotated = annotate_candidates(
            image=base_image,
            candidates=candidates,
            slide_width=400,
            slide_height=300,
        )

        self.assertEqual(annotated.getpixel((100, 80)), (255, 59, 48))

    def test_label_layout_scales_with_box_size(self) -> None:
        image = Image.new("RGB", (600, 400), "white")
        draw = ImageDraw.Draw(image)

        small_layout = _fit_label_layout(
            draw=draw,
            text="12",
            box=(20, 20, 90, 60),
        )
        large_layout = _fit_label_layout(
            draw=draw,
            text="12",
            box=(20, 20, 260, 180),
        )

        self.assertGreater(large_layout["label_height"], small_layout["label_height"])
        self.assertGreater(large_layout["label_width"], small_layout["label_width"])

    def test_show_command_writes_json_and_image(self) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)
            pptx_path = tmp_dir / "demo.pptx"
            out_path = tmp_dir / "out.png"
            candidates_path = tmp_dir / "candidates.json"

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            text_box.text_frame.text = "Quarterly Review"
            presentation.save(pptx_path)

            stdout = StringIO()
            with patch("pptx_cli.show.render_slide_preview", return_value=Image.new("RGB", (800, 600), "white")):
                with redirect_stdout(stdout):
                    exit_code = main(
                        [
                            "show",
                            "--input",
                            str(pptx_path),
                            "--slide",
                            "0",
                            "--annotate",
                            "--output",
                            str(out_path),
                            "--candidates-out",
                            str(candidates_path),
                        ]
                    )

            self.assertEqual(exit_code, 0)
            self.assertTrue(out_path.exists())
            self.assertTrue(candidates_path.exists())
            payload = json.loads(stdout.getvalue())
            self.assertEqual(payload["command"], "show")
            self.assertTrue(payload["annotated"])
            self.assertEqual(len(payload["candidates"]), 1)


if __name__ == "__main__":
    unittest.main()
