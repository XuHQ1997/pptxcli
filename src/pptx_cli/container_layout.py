from __future__ import annotations

import math
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from .models import BBox

LAYOUT_TYPES = {"vertical", "horizontal", "grid"}
LEAF_TYPES = {"text", "image", "svg"}
IMAGE_FITS = {"fill", "contain", "cover"}
LEGACY_FIELDS = {
    "bbox",
    "x",
    "y",
    "w",
    "h",
    "width",
    "height",
    "min_width",
    "min_height",
    "max_width",
    "max_height",
    "align",
    "cross_align",
    "padding",
    "rows",
    "row",
    "column",
    "row_span",
    "col_span",
    "value",
    "content",
    "src",
}
TEXT_STYLE_FIELDS = {"font_name", "font_size", "bold", "italic", "underline", "color"}


@dataclass(slots=True)
class ContentNode:
    node_type: str
    ratio_x: float = 1.0
    ratio_y: float = 1.0
    layout: str | None = None
    gap: float = 0.0
    columns: int | None = None
    text: str | None = None
    path: str | None = None
    fit: str = "contain"
    name: str | None = None
    style: dict[str, Any] = field(default_factory=dict)
    children: list["ContentNode"] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        payload: dict[str, Any] = {
            "ratio": _serialize_ratio(self.ratio_x, self.ratio_y),
            "name": self.name,
        }
        if self.layout is not None:
            payload["layout"] = self.layout
            payload["gap"] = self.gap
            payload["columns"] = self.columns
            payload["children"] = [child.to_dict() for child in self.children]
        else:
            payload["type"] = self.node_type
            payload["fit"] = self.fit
            payload["style"] = self.style
            if self.text is not None:
                payload["text"] = self.text
            if self.path is not None:
                payload["path"] = self.path
        return payload


@dataclass(slots=True)
class ResolvedContent:
    node_type: str
    bbox: BBox
    name: str | None = None
    text: str | None = None
    path: str | None = None
    style: dict[str, Any] = field(default_factory=dict)
    fit: str = "contain"
    children: list["ResolvedContent"] = field(default_factory=list)
    debug: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        payload: dict[str, Any] = {
            "type": self.node_type,
            "bbox": self.bbox.to_dict(),
            "name": self.name,
            "fit": self.fit,
            "style": self.style,
            "debug": self.debug,
            "children": [child.to_dict() for child in self.children],
        }
        if self.text is not None:
            payload["text"] = self.text
        if self.path is not None:
            payload["path"] = self.path
        return payload


def parse_content_spec(
    raw_value: str,
    *,
    slide_width: int,
    slide_height: int,
    default_bbox: BBox | None = None,
) -> ContentNode:
    try:
        payload = json_loads(raw_value)
    except ValueError as exc:
        raise ValueError(f"invalid --content JSON: {exc}") from exc
    if not isinstance(payload, dict):
        raise ValueError("--content expects a JSON object")
    _ = default_bbox or BBox(0, 0, slide_width, slide_height)
    return _parse_node(payload, is_root=True)


def solve_content_layout(
    node: ContentNode,
    *,
    slide_width: int,
    slide_height: int,
    default_bbox: BBox | None = None,
) -> ResolvedContent:
    root_bbox = default_bbox or BBox(0, 0, slide_width, slide_height)
    bounded_root = _clip_bbox(root_bbox, BBox(0, 0, slide_width, slide_height))
    return _resolve_node(node=node, bbox=bounded_root)


def flatten_resolved_leaves(node: ResolvedContent) -> list[ResolvedContent]:
    if node.node_type in LEAF_TYPES:
        return [node]
    leaves: list[ResolvedContent] = []
    for child in node.children:
        leaves.extend(flatten_resolved_leaves(child))
    return leaves


def render_resolved_content(*, slide: Any, resolved: ResolvedContent) -> list[dict[str, Any]]:
    rendered: list[dict[str, Any]] = []
    for leaf in flatten_resolved_leaves(resolved):
        if leaf.node_type == "text":
            shape = _render_text(slide=slide, leaf=leaf)
        elif leaf.node_type in {"image", "svg"}:
            shape = _render_image(slide=slide, leaf=leaf)
        else:
            raise ValueError(f"unsupported resolved content type: {leaf.node_type}")
        rendered.append(
            {
                "type": leaf.node_type,
                "shape_id": int(shape.shape_id),
                "bbox": leaf.bbox.to_dict(),
                "name": leaf.name,
            }
        )
    return rendered


def json_loads(raw_value: str) -> Any:
    import json

    return json.loads(raw_value)


def _parse_node(payload: dict[str, Any], *, is_root: bool) -> ContentNode:
    _reject_legacy_fields(payload)
    layout = _read_optional_str(payload, "layout")
    if layout is not None:
        if layout not in LAYOUT_TYPES:
            raise ValueError(f"unsupported layout type: {layout}")
        if "type" in payload:
            raise ValueError("container node should use `layout` only; `type=container` is no longer supported")
        children = payload.get("children")
        if not isinstance(children, list):
            raise ValueError("container node requires children array")
        return ContentNode(
            node_type="container",
            layout=layout,
            ratio_x=_read_ratio(payload, "ratio", axis="x", default=1.0),
            ratio_y=_read_ratio(payload, "ratio", axis="y", default=1.0),
            gap=_read_non_negative_number(payload, "gap", default=0.0),
            columns=_read_optional_int(payload, "columns"),
            name=_read_optional_str(payload, "name"),
            children=[
                _parse_node(child, is_root=False) if isinstance(child, dict) else _invalid_child(child)
                for child in children
            ],
        )

    node_type = _read_optional_str(payload, "type")
    if node_type not in LEAF_TYPES:
        if is_root:
            raise ValueError("root content node must use `layout`")
        raise ValueError(f"unsupported content node type: {node_type or '<empty>'}")
    node = ContentNode(
        node_type=node_type,
        ratio_x=_read_ratio(payload, "ratio", axis="x", default=1.0),
        ratio_y=_read_ratio(payload, "ratio", axis="y", default=1.0),
        fit=_read_fit(payload.get("fit")),
        name=_read_optional_str(payload, "name"),
        style=_read_style(payload.get("style")),
    )
    if node_type == "text":
        node.text = _read_required_str(payload, "text")
    else:
        node.path = _resolve_asset_path(_read_required_str(payload, "path"))
    return node


def _invalid_child(value: Any) -> ContentNode:
    raise ValueError(f"container child must be an object, got: {type(value).__name__}")


def _resolve_node(*, node: ContentNode, bbox: BBox) -> ResolvedContent:
    if node.layout is None:
        return ResolvedContent(
            node_type=node.node_type,
            bbox=bbox,
            name=node.name,
            text=node.text,
            path=node.path,
            style=node.style,
            fit=node.fit,
            debug={"ratio": _serialize_ratio(node.ratio_x, node.ratio_y)},
        )

    if node.layout in {"vertical", "horizontal"}:
        children = _resolve_linear_children(node=node, inner=bbox)
    elif node.layout == "grid":
        children = _resolve_grid_children(node=node, inner=bbox)
    else:
        raise ValueError(f"unsupported container layout: {node.layout}")
    return ResolvedContent(
        node_type="container",
        bbox=bbox,
        name=node.name,
        children=children,
        debug={"layout": node.layout, "gap": node.gap},
    )


def _resolve_linear_children(*, node: ContentNode, inner: BBox) -> list[ResolvedContent]:
    is_vertical = node.layout == "vertical"
    child_count = len(node.children)
    if child_count == 0:
        return []
    main_axis = inner.h if is_vertical else inner.w
    gap_size = int(round(main_axis * node.gap))
    gap_total = gap_size * max(child_count - 1, 0)
    available_main = max(main_axis - gap_total, 0)
    ratios = [child.ratio_y if is_vertical else child.ratio_x for child in node.children]
    sizes = _distribute_by_ratio(available_main, ratios)
    resolved_children: list[ResolvedContent] = []
    cursor = 0
    for child, main_size in zip(node.children, sizes, strict=False):
        if is_vertical:
            child_bbox = BBox(inner.x, inner.y + cursor, inner.w, main_size)
        else:
            child_bbox = BBox(inner.x + cursor, inner.y, main_size, inner.h)
        resolved_children.append(_resolve_node(node=child, bbox=child_bbox))
        cursor += main_size + gap_size
    return resolved_children


def _resolve_grid_children(*, node: ContentNode, inner: BBox) -> list[ResolvedContent]:
    child_count = len(node.children)
    if child_count == 0:
        return []
    columns = node.columns or max(int(math.ceil(math.sqrt(child_count))), 1)
    rows = max(int(math.ceil(child_count / columns)), 1)
    column_gap = int(round(inner.w * node.gap))
    row_gap = int(round(inner.h * node.gap))
    column_ratios, row_ratios = _derive_grid_track_ratios(node.children, columns=columns, rows=rows)
    column_widths = _distribute_by_ratio(max(inner.w - column_gap * max(columns - 1, 0), 0), column_ratios)
    row_heights = _distribute_by_ratio(max(inner.h - row_gap * max(rows - 1, 0), 0), row_ratios)
    column_positions = _build_track_positions(inner.x, column_widths, column_gap)
    row_positions = _build_track_positions(inner.y, row_heights, row_gap)
    resolved_children: list[ResolvedContent] = []
    for index, child in enumerate(node.children):
        row = index // columns
        column = index % columns
        child_bbox = BBox(
            x=column_positions[column],
            y=row_positions[row],
            w=column_widths[column],
            h=row_heights[row],
        )
        resolved_children.append(_resolve_node(node=child, bbox=child_bbox))
    return resolved_children


def _clip_bbox(bbox: BBox, bounds: BBox) -> BBox:
    x = min(max(bbox.x, bounds.x), bounds.x + bounds.w)
    y = min(max(bbox.y, bounds.y), bounds.y + bounds.h)
    max_w = max(bounds.x + bounds.w - x, 0)
    max_h = max(bounds.y + bounds.h - y, 0)
    return BBox(x=x, y=y, w=min(max(bbox.w, 0), max_w), h=min(max(bbox.h, 0), max_h))


def _distribute_by_ratio(total: int, ratios: list[float]) -> list[int]:
    ratio_total = sum(ratios)
    if ratio_total <= 0:
        ratios = [1.0] * len(ratios)
        ratio_total = float(len(ratios))
    sizes: list[int] = []
    consumed = 0
    for index, ratio in enumerate(ratios):
        if index == len(ratios) - 1:
            size = max(total - consumed, 0)
        else:
            size = int(round(total * (ratio / ratio_total)))
            size = min(size, max(total - consumed, 0))
        consumed += size
        sizes.append(size)
    return sizes


def _derive_grid_track_ratios(
    children: list[ContentNode],
    *,
    columns: int,
    rows: int,
) -> tuple[list[float], list[float]]:
    column_ratios = [1.0] * columns
    row_ratios = [1.0] * rows
    for index, child in enumerate(children):
        row = index // columns
        column = index % columns
        column_ratios[column] = max(column_ratios[column], child.ratio_x)
        row_ratios[row] = max(row_ratios[row], child.ratio_y)
    return column_ratios, row_ratios


def _build_track_positions(start: int, sizes: list[int], gap: int) -> list[int]:
    positions: list[int] = []
    cursor = start
    for size in sizes:
        positions.append(cursor)
        cursor += size + gap
    return positions


def _render_text(*, slide: Any, leaf: ResolvedContent) -> Any:
    shape = slide.shapes.add_textbox(leaf.bbox.x, leaf.bbox.y, leaf.bbox.w, leaf.bbox.h)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = (leaf.text or "").split("\n")
    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.clear()
    for index, line in enumerate(lines):
        paragraph = first_paragraph if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = line
        _apply_run_style(run=run, style=leaf.style)
    return shape


def _render_image(*, slide: Any, leaf: ResolvedContent) -> Any:
    if leaf.path is None:
        raise ValueError("image/svg leaf missing path")
    path = Path(leaf.path)
    left, top, width, height = _resolve_image_box(path=path, bbox=leaf.bbox, fit=leaf.fit)
    return slide.shapes.add_picture(str(path), left, top, width, height)


def _resolve_image_box(*, path: Path, bbox: BBox, fit: str) -> tuple[int, int, int, int]:
    if fit not in IMAGE_FITS:
        fit = "contain"
    if fit == "fill" or path.suffix.lower() == ".svg":
        return bbox.x, bbox.y, bbox.w, bbox.h
    with Image.open(path) as image:
        image_width, image_height = image.size
    if image_width <= 0 or image_height <= 0 or bbox.w <= 0 or bbox.h <= 0:
        return bbox.x, bbox.y, bbox.w, bbox.h
    image_ratio = image_width / image_height
    box_ratio = bbox.w / bbox.h
    if (fit == "contain" and image_ratio >= box_ratio) or (fit == "cover" and image_ratio <= box_ratio):
        width = bbox.w
        height = int(round(width / image_ratio))
    else:
        height = bbox.h
        width = int(round(height * image_ratio))
    left = bbox.x + max((bbox.w - width) // 2, 0)
    top = bbox.y + max((bbox.h - height) // 2, 0)
    return left, top, width, height


def _apply_run_style(*, run: Any, style: dict[str, Any]) -> None:
    font = run.font
    if "font_name" in style:
        font.name = str(style["font_name"])
    if "font_size" in style:
        font.size = Pt(float(style["font_size"]))
    if "bold" in style:
        font.bold = bool(style["bold"])
    if "italic" in style:
        font.italic = bool(style["italic"])
    if "underline" in style:
        font.underline = bool(style["underline"])
    color = style.get("color")
    rgb = _parse_color(color)
    if rgb is not None:
        font.color.rgb = rgb


def _parse_color(value: Any) -> RGBColor | None:
    if value is None:
        return None
    raw = str(value).strip().lstrip("#")
    if len(raw) != 6:
        raise ValueError(f"invalid hex color: {value}")
    try:
        return RGBColor.from_string(raw.upper())
    except ValueError as exc:
        raise ValueError(f"invalid hex color: {value}") from exc


def _resolve_asset_path(raw_value: str) -> str:
    candidate = Path(raw_value).expanduser().resolve()
    if not candidate.exists():
        raise ValueError(f"content asset does not exist: {candidate}")
    return str(candidate)


def _read_style(value: Any) -> dict[str, Any]:
    if value is None:
        return {}
    if not isinstance(value, dict):
        raise ValueError("content style must be an object")
    unsupported = [key for key in value if key not in TEXT_STYLE_FIELDS]
    if unsupported:
        raise ValueError(f"unsupported text style fields: {', '.join(sorted(unsupported))}")
    return dict(value)


def _read_ratio(
    payload: dict[str, Any],
    key: str,
    *,
    axis: str,
    default: float,
) -> float:
    value = payload.get(key, default)
    if isinstance(value, list):
        if len(value) != 2:
            raise ValueError(f"{key} array must contain exactly two numbers")
        ratios = [_coerce_positive_number(item, key=key) for item in value]
        return ratios[0] if axis == "x" else ratios[1]
    scalar = _coerce_positive_number(value, key=key)
    return scalar


def _coerce_positive_number(value: Any, *, key: str) -> float:
    if isinstance(value, bool):
        raise ValueError(f"{key} must be a number or [number, number]")
    if isinstance(value, (int, float)):
        number = float(value)
    elif isinstance(value, str) and value.strip():
        try:
            number = float(value)
        except ValueError as exc:
            raise ValueError(f"{key} must be a number or [number, number]") from exc
    else:
        raise ValueError(f"{key} must be a number or [number, number]")
    if number <= 0:
        raise ValueError(f"{key} must be > 0")
    return number


def _serialize_ratio(ratio_x: float, ratio_y: float) -> float | list[float]:
    if math.isclose(ratio_x, ratio_y):
        return ratio_x
    return [ratio_x, ratio_y]


def _read_fit(value: Any) -> str:
    if value is None:
        return "contain"
    fit = str(value).strip().lower()
    if fit not in IMAGE_FITS:
        raise ValueError(f"unsupported fit value: {value}")
    return fit


def _read_non_negative_number(payload: dict[str, Any], key: str, default: float) -> float:
    value = payload.get(key, default)
    if isinstance(value, bool):
        raise ValueError(f"{key} must be a number")
    if isinstance(value, (int, float)):
        number = float(value)
    elif isinstance(value, str) and value.strip():
        try:
            number = float(value)
        except ValueError as exc:
            raise ValueError(f"{key} must be a number") from exc
    else:
        raise ValueError(f"{key} must be a number")
    if number < 0:
        raise ValueError(f"{key} must be >= 0")
    return number


def _read_optional_str(payload: dict[str, Any], key: str) -> str | None:
    value = payload.get(key)
    if value is None:
        return None
    text = str(value).strip()
    return text or None


def _read_required_str(payload: dict[str, Any], key: str) -> str:
    value = _read_optional_str(payload, key)
    if value is None:
        raise ValueError(f"missing required field: {key}")
    return value


def _read_optional_int(payload: dict[str, Any], key: str) -> int | None:
    value = payload.get(key)
    if value is None:
        return None
    if isinstance(value, bool):
        raise ValueError(f"{key} must be an integer")
    if isinstance(value, (int, float)):
        return int(value)
    if isinstance(value, str) and value.strip():
        try:
            return int(float(value))
        except ValueError as exc:
            raise ValueError(f"{key} must be an integer") from exc
    raise ValueError(f"{key} must be an integer")


def _read_int(payload: dict[str, Any], key: str, default: int = 0) -> int:
    value = payload.get(key, default)
    result = _read_optional_int({key: value}, key)
    if result is None:
        return default
    return result


def _reject_legacy_fields(payload: dict[str, Any]) -> None:
    legacy = [key for key in payload if key in LEGACY_FIELDS]
    if legacy:
        raise ValueError(
            "unsupported legacy container fields: "
            + ", ".join(sorted(legacy))
            + ". Use ratio-based layout instead."
        )
