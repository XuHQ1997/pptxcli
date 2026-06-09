from __future__ import annotations

import io
import json
import os
import re
import shutil
import unicodedata
from copy import deepcopy
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from .container_layout import (
    parse_content_spec,
    render_resolved_content,
    solve_content_layout,
)
from .session import (
    SessionError,
    cleanup_stale_session_state,
    ensure_session_for_origin,
    load_session_state_if_exists,
    load_session_state,
    resolve_state_file_path,
    session_request,
    update_session_mode,
)

TEMPLATE_ROOT_ENV = "PPTXCLI_TEMPLATE_ROOT"
TEMPLATE_DRAFT_VERSION = 1
MANIFEST_VERSION = 1
EDIT_WORKING_FILE_NAME = ".pptxcli-editing.pptx"


def resolve_template_root(cli_argv0: str | None = None) -> Path:
    env_path = os.environ.get(TEMPLATE_ROOT_ENV)
    if env_path:
        return Path(env_path).expanduser().resolve()
    return resolve_state_file_path(cli_argv0).parent / "templates"


def create_template_draft(
    *,
    name: str,
    from_file: Path | None = None,
) -> dict[str, Any]:
    template_name = normalize_template_name(name)
    origin_file = _resolve_origin_file(from_file=from_file)
    draft_path = resolve_template_root() / f"{template_name}.json"
    if draft_path.exists():
        raise ValueError(f"template draft already exists: {draft_path}")

    payload = {
        "draft_version": TEMPLATE_DRAFT_VERSION,
        "template_name": template_name,
        "origin_file": str(origin_file),
        "created_at": _now_iso(),
        "slides": [],
    }
    _write_json(draft_path, payload)
    _persist_workflow_state(
        mode="template_extract",
        active_template=template_name,
        edit_context=None,
    )
    return {
        "command": "template create",
        "status": "created",
        "template_name": template_name,
        "origin_file": str(origin_file),
        "draft_path": str(draft_path),
    }


def add_slide_to_template_draft(
    *,
    slide_index: int,
    field_specs: list[str],
    replace: bool,
) -> dict[str, Any]:
    template_name = _resolve_template_name()
    draft_path = resolve_template_root() / f"{template_name}.json"
    draft = _load_template_draft(draft_path)
    inspect_payload = _load_slide_candidates(slide_index=slide_index)
    slide_record = _build_slide_record(
        slide_index=slide_index,
        field_specs=field_specs,
        inspect_payload=inspect_payload,
    )

    existing_index = _find_slide_record_index(
        draft=draft,
        source_slide_index=slide_index,
    )
    if existing_index is not None:
        if not replace:
            raise ValueError(
                "slide already exists in template draft; pass --replace to overwrite it"
            )
        draft["slides"][existing_index] = slide_record
    else:
        draft["slides"].append(slide_record)

    draft["updated_at"] = _now_iso()
    _write_json(draft_path, draft)
    _persist_workflow_state(
        mode="template_extract",
        active_template=template_name,
        edit_context=None,
    )
    return {
        "command": "template add_slide",
        "status": "updated",
        "template_name": template_name,
        "draft_path": str(draft_path),
        "slide_name": slide_record["slide_name"],
        "source_slide_index": slide_index,
        "field_count": len(slide_record["fields"]),
    }


def save_template_package(
) -> dict[str, Any]:
    template_name = _resolve_template_name()
    draft_path = resolve_template_root() / f"{template_name}.json"
    draft = _load_template_draft(draft_path)
    slides = draft.get("slides", [])
    if not slides:
        raise ValueError("template draft has no slides; add at least one slide before save")

    origin_file = Path(str(draft["origin_file"])).resolve()
    if not origin_file.exists():
        raise ValueError(f"origin file does not exist: {origin_file}")

    package_dir = resolve_template_root() / template_name
    package_dir.mkdir(parents=True, exist_ok=True)
    template_pptx_path = package_dir / "template.pptx"
    manifest_path = package_dir / "manifest.json"

    source_slide_indexes = [int(slide["source_slide_index"]) for slide in slides]
    _save_cropped_presentation(
        origin_file=origin_file,
        slide_indexes=source_slide_indexes,
        output_path=template_pptx_path,
    )

    manifest = _build_manifest(
        draft=draft,
        draft_path=draft_path,
        template_pptx_path=template_pptx_path,
    )
    _write_json(manifest_path, manifest)
    _persist_workflow_state(
        mode="template_extract",
        active_template=template_name,
        edit_context=None,
    )

    return {
        "command": "template save",
        "status": "saved",
        "template_name": template_name,
        "draft_path": str(draft_path),
        "package_dir": str(package_dir),
        "template_pptx_path": str(template_pptx_path),
        "manifest_path": str(manifest_path),
        "slide_count": manifest["slide_count"],
        "field_count": manifest["field_count"],
    }


def build_template_package(
) -> dict[str, Any]:
    payload = save_template_package()
    payload["command"] = "template build"
    return payload


def create_edit_presentation(
    *,
    output_path: Path,
    template_ref: str | None = None,
) -> dict[str, Any]:
    state = _load_session_metadata()
    edit_context = state.get("edit_context")
    if isinstance(edit_context, dict) and str(edit_context.get("working_pptx_path", "")).strip():
        raise ValueError("an edit draft is already active; run `pptxcli edit save` or wait for timeout first")

    manifest_path, template_pptx_path = _resolve_template_package_paths(
        template_ref=template_ref,
    )
    manifest = _load_manifest(manifest_path)
    working_path = resolve_state_file_path().parent / EDIT_WORKING_FILE_NAME

    output_presentation = Presentation(str(template_pptx_path))
    _remove_all_slides(output_presentation)
    working_path.parent.mkdir(parents=True, exist_ok=True)
    output_presentation.save(str(working_path))

    template_name = str(manifest.get("template_name") or _derive_template_name(manifest_path))
    context = {
        "template_name": template_name,
        "manifest_path": str(manifest_path),
        "template_pptx_path": str(template_pptx_path),
        "output_path": str(output_path.resolve()),
        "working_pptx_path": str(working_path),
        "slide_count": 0,
        "filled_slides": [],
        "created_at": _now_iso(),
    }
    _persist_workflow_state(
        mode="edit_ppt",
        active_template=template_name,
        edit_context=context,
    )
    return {
        "command": "edit create",
        "status": "created",
        "template_name": template_name,
        "manifest_path": str(manifest_path),
        "template_pptx_path": str(template_pptx_path),
        "working_pptx_path": str(working_path),
        "output_path": str(output_path.resolve()),
        "slide_count": 0,
    }


def show_template_fields_for_edit(
    *,
    slide_index: int,
) -> dict[str, Any]:
    state = _load_session_metadata(required=True)
    manifest_path, _template_pptx_path = _resolve_edit_template_package_paths(state)
    manifest = _load_manifest(manifest_path)
    manifest_slide = _resolve_manifest_slide(
        manifest=manifest,
        slide_index=slide_index,
    )
    template_name = str(
        state.get("active_template")
        or manifest.get("template_name")
        or _derive_template_name(manifest_path)
    ).strip()
    return {
        "command": "edit show_template",
        "status": "ok",
        "mode": state.get("mode", "idle"),
        "template_name": template_name,
        "manifest_path": str(manifest_path),
        "slide": {
            "slide_name": str(manifest_slide["slide_name"]),
            "template_slide_index": int(manifest_slide["template_slide_index"]),
            "source_slide_index": manifest_slide.get("source_slide_index"),
            "slide_size": manifest_slide.get("slide_size"),
            "field_count": len(manifest_slide["fields"]),
            "fields": [_serialize_template_field(field) for field in manifest_slide["fields"]],
        },
    }


def fill_template_into_edit_presentation(
    *,
    slide_index: int,
    field_specs: list[str] | None,
    content_spec: str | None = None,
) -> dict[str, Any]:
    state = _load_session_metadata(required=True)
    edit_context = _require_edit_context(state)
    manifest_path = Path(str(edit_context["manifest_path"])).resolve()
    template_pptx_path = Path(str(edit_context["template_pptx_path"])).resolve()
    working_pptx_path = Path(str(edit_context["working_pptx_path"])).resolve()

    manifest = _load_manifest(manifest_path)
    requested_slide = _build_fill_request(
        slide_index=slide_index,
        field_specs=field_specs,
        manifest=manifest,
    )
    if not working_pptx_path.exists():
        raise ValueError(f"edit working file does not exist: {working_pptx_path}")

    source_presentation = Presentation(str(template_pptx_path))
    output_presentation = Presentation(str(working_pptx_path))
    template_slide = source_presentation.slides[requested_slide["template_slide_index"]]
    output_slide, shape_map = _clone_slide(
        source_slide=template_slide,
        destination_presentation=output_presentation,
    )
    _apply_slide_fields(
        slide=output_slide,
        shape_map=shape_map,
        assignments=requested_slide["assignments"],
    )
    rendered_content: list[dict[str, Any]] = []
    resolved_content_payload: dict[str, Any] | None = None
    if content_spec is not None:
        content_tree = parse_content_spec(
            content_spec,
            slide_width=int(output_presentation.slide_width),
            slide_height=int(output_presentation.slide_height),
        )
        resolved_content = solve_content_layout(
            content_tree,
            slide_width=int(output_presentation.slide_width),
            slide_height=int(output_presentation.slide_height),
        )
        rendered_content = render_resolved_content(
            slide=output_slide,
            resolved=resolved_content,
        )
        resolved_content_payload = resolved_content.to_dict()
    output_presentation.save(str(working_pptx_path))

    filled_slides = list(edit_context.get("filled_slides") or [])
    filled_slides.append(
        {
            "slide_name": requested_slide["slide_name"],
            "template_slide_index": requested_slide["template_slide_index"],
            "field_count": len(requested_slide["assignments"]),
        }
    )
    updated_context = {
        **edit_context,
        "slide_count": len(output_presentation.slides),
        "filled_slides": filled_slides,
        "updated_at": _now_iso(),
    }
    _persist_workflow_state(
        mode="edit_ppt",
        active_template=str(state.get("active_template") or edit_context.get("template_name") or ""),
        edit_context=updated_context,
    )
    return {
        "command": "edit fill_template",
        "status": "updated",
        "output_path": str(Path(str(updated_context["output_path"])).resolve()),
        "working_pptx_path": str(working_pptx_path),
        "slide_count": len(output_presentation.slides),
        "field_count": len(requested_slide["assignments"]),
        "content_count": len(rendered_content),
        "slide": requested_slide["slide_name"],
        "template_slide_index": requested_slide["template_slide_index"],
        "content_layout": resolved_content_payload,
        "rendered_content": rendered_content,
    }


def save_edit_presentation() -> dict[str, Any]:
    state = _load_session_metadata(required=True)
    edit_context = _require_edit_context(state)
    working_pptx_path = Path(str(edit_context["working_pptx_path"])).resolve()
    output_path = Path(str(edit_context["output_path"])).resolve()
    if not working_pptx_path.exists():
        raise ValueError(f"edit working file does not exist: {working_pptx_path}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(working_pptx_path, output_path)
    working_pptx_path.unlink()

    active_template = str(state.get("active_template", "")).strip()
    next_mode = "template_extract" if active_template else "idle"
    _persist_workflow_state(
        mode=next_mode,
        active_template=active_template or None,
        edit_context=None,
    )
    return {
        "command": "edit save",
        "status": "saved",
        "output_path": str(output_path),
        "slide_count": int(edit_context.get("slide_count", 0)),
        "template_name": edit_context.get("template_name"),
    }


def normalize_template_name(name: str) -> str:
    normalized = Path(name.strip()).name
    if normalized.endswith(".json") or normalized.endswith(".pptx"):
        normalized = Path(normalized).stem
    if not normalized:
        raise ValueError("template name cannot be empty")
    return normalized


def _resolve_template_name() -> str:
    state = _load_session_metadata(required=True)
    active_template = str(state.get("active_template", "")).strip()
    if not active_template:
        raise ValueError(
            "no active template found. Run `pptxcli template create --name <name>` first."
        )
    return normalize_template_name(active_template)


def _resolve_origin_file(*, from_file: Path | None = None) -> Path:
    state_file = resolve_state_file_path()
    if from_file is not None:
        origin_file = from_file.expanduser().resolve()
        try:
            ensure_session_for_origin(origin_file=origin_file, state_file=state_file)
        except SessionError as exc:
            cleanup_stale_session_state(state_file)
            raise ValueError(str(exc)) from exc
        return origin_file

    try:
        session_request(state_file=state_file, method="GET", route="/health", timeout=1.0)
        state = load_session_state(state_file)
    except SessionError as exc:
        cleanup_stale_session_state(state_file)
        raise ValueError(str(exc)) from exc

    origin_file = Path(str(state.get("origin_file", ""))).resolve()
    if not origin_file.exists():
        raise ValueError(f"origin file does not exist: {origin_file}")
    return origin_file


def _load_slide_candidates(*, slide_index: int) -> dict[str, Any]:
    state_file = resolve_state_file_path()
    try:
        return session_request(
            state_file=state_file,
            method="POST",
            route="/template/candidates",
            payload={
                "command_name": "template add_slide",
                "slide_index": slide_index,
            },
        )
    except SessionError as exc:
        cleanup_stale_session_state(state_file)
        raise ValueError(str(exc)) from exc


def _build_slide_record(
    *,
    slide_index: int,
    field_specs: list[str],
    inspect_payload: dict[str, Any],
) -> dict[str, Any]:
    candidates = inspect_payload.get("candidates")
    if not isinstance(candidates, list):
        raise ValueError("invalid inspect payload: missing candidates")
    candidates_by_index = {
        int(candidate.get("index")): candidate
        for candidate in candidates
        if isinstance(candidate, dict) and candidate.get("index") is not None
    }

    if not field_specs:
        raise ValueError("at least one --field is required")

    seen_indexes: set[int] = set()
    seen_field_keys: set[str] = set()
    fields: list[dict[str, Any]] = []
    for raw_field_spec in field_specs:
        selected_index, description = _parse_field_spec(raw_field_spec)
        if selected_index in seen_indexes:
            raise ValueError(f"duplicate field index on slide {slide_index}: {selected_index}")
        seen_indexes.add(selected_index)

        candidate = candidates_by_index.get(selected_index)
        if candidate is None:
            raise ValueError(
                f"field index {selected_index} does not exist on slide {slide_index}"
            )

        field_type = str(candidate.get("kind", "")).strip()
        field_key = _build_field_key(description=description, used_keys=seen_field_keys)
        field_record: dict[str, Any] = {
            "index": selected_index,
            "field_key": field_key,
            "description": description,
            "type": field_type,
            "shape_id": int(candidate["shape_id"]),
            "role_hint": candidate.get("role_hint"),
            "bbox": candidate.get("bbox"),
        }
        if candidate.get("text_excerpt") is not None:
            field_record["text_excerpt"] = candidate["text_excerpt"]
        if candidate.get("image_name") is not None:
            field_record["image_name"] = candidate["image_name"]
        fields.append(field_record)

    return {
        "slide_name": f"slide_{slide_index}",
        "source_slide_index": slide_index,
        "slide_size": inspect_payload.get("slide_size"),
        "fields": fields,
        "added_at": _now_iso(),
    }


def _find_slide_record_index(
    *,
    draft: dict[str, Any],
    source_slide_index: int,
) -> int | None:
    for index, slide in enumerate(draft.get("slides", [])):
        existing_source_slide = int(slide.get("source_slide_index", -1))
        if existing_source_slide == source_slide_index:
            return index
    return None


def _build_manifest(
    *,
    draft: dict[str, Any],
    draft_path: Path,
    template_pptx_path: Path,
) -> dict[str, Any]:
    slides_payload: list[dict[str, Any]] = []
    all_fields: list[dict[str, Any]] = []

    for template_slide_index, slide in enumerate(draft["slides"]):
        fields_payload: list[dict[str, Any]] = []
        for field in slide["fields"]:
            field_payload = {
                "index": field["index"],
                "field_key": field["field_key"],
                "description": field["description"],
                "type": field["type"],
                "shape_id": field["shape_id"],
                "role_hint": field.get("role_hint"),
                "bbox": field.get("bbox"),
                "binding": {
                    "kind": "shape",
                    "shape_id": field["shape_id"],
                },
            }
            if field.get("text_excerpt") is not None:
                field_payload["text_excerpt"] = field["text_excerpt"]
            if field.get("image_name") is not None:
                field_payload["image_name"] = field["image_name"]
            fields_payload.append(field_payload)
            all_fields.append(
                {
                    "slide_name": slide["slide_name"],
                    **field_payload,
                }
            )

        slides_payload.append(
            {
                "slide_name": slide["slide_name"],
                "template_slide_index": template_slide_index,
                "source_slide_index": slide["source_slide_index"],
                "slide_size": slide.get("slide_size"),
                "field_count": len(fields_payload),
                "fields": fields_payload,
            }
        )

    return {
        "manifest_version": MANIFEST_VERSION,
        "template_name": draft["template_name"],
        "origin_file": draft["origin_file"],
        "draft_path": str(draft_path),
        "template_file": template_pptx_path.name,
        "created_at": _now_iso(),
        "slide_count": len(slides_payload),
        "field_count": len(all_fields),
        "slides": slides_payload,
        "fields": all_fields,
    }


def _save_cropped_presentation(
    *,
    origin_file: Path,
    slide_indexes: list[int],
    output_path: Path,
) -> None:
    if len(set(slide_indexes)) != len(slide_indexes):
        raise ValueError("template draft contains duplicate source_slide_index values")

    presentation = Presentation(str(origin_file))
    total_slides = len(presentation.slides)
    for slide_index in slide_indexes:
        if slide_index < 0 or slide_index >= total_slides:
            raise ValueError(
                f"slide index {slide_index} out of range; total slides: {total_slides}"
            )

    slide_id_list = presentation.slides._sldIdLst
    original_slide_ids = list(slide_id_list)
    selected_by_index = {slide_index: original_slide_ids[slide_index] for slide_index in slide_indexes}

    for slide_index in range(len(original_slide_ids) - 1, -1, -1):
        slide_id = original_slide_ids[slide_index]
        if slide_index in selected_by_index:
            continue
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)

    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_index in slide_indexes:
        slide_id_list.append(selected_by_index[slide_index])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))


def _resolve_template_package_paths(
    *,
    template_ref: str | None,
) -> tuple[Path, Path]:
    candidate_ref = template_ref

    if candidate_ref:
        named_template_dir = resolve_template_root() / normalize_template_name(candidate_ref)
        if named_template_dir.is_dir():
            manifest_path = named_template_dir / "manifest.json"
            manifest = _load_manifest(manifest_path)
            template_pptx_path = named_template_dir / str(manifest.get("template_file", "template.pptx"))
            if not template_pptx_path.exists():
                raise ValueError(f"template file does not exist: {template_pptx_path}")
            return manifest_path, template_pptx_path

        reference_path = Path(candidate_ref).expanduser().resolve()
        if reference_path.is_dir():
            manifest_path = reference_path / "manifest.json"
            manifest = _load_manifest(manifest_path)
            template_pptx_path = reference_path / str(manifest.get("template_file", "template.pptx"))
            if not template_pptx_path.exists():
                raise ValueError(f"template file does not exist: {template_pptx_path}")
            return manifest_path, template_pptx_path

        if reference_path.name == "manifest.json":
            manifest = _load_manifest(reference_path)
            template_pptx_path = reference_path.parent / str(
                manifest.get("template_file", "template.pptx")
            )
            if not template_pptx_path.exists():
                raise ValueError(f"template file does not exist: {template_pptx_path}")
            return reference_path, template_pptx_path

        if reference_path.suffix.lower() == ".pptx":
            manifest_path = reference_path.parent / "manifest.json"
            _load_manifest(manifest_path)
            return manifest_path, reference_path

        raise ValueError(
            "template reference must be a template directory, manifest.json, or template.pptx"
        )

    template_name = _resolve_template_name()
    package_dir = resolve_template_root() / template_name
    manifest_path = package_dir / "manifest.json"
    manifest = _load_manifest(manifest_path)
    template_pptx_path = package_dir / str(manifest.get("template_file", "template.pptx"))
    if not template_pptx_path.exists():
        raise ValueError(f"template file does not exist: {template_pptx_path}")
    return manifest_path, template_pptx_path


def _resolve_edit_template_package_paths(state: dict[str, Any]) -> tuple[Path, Path]:
    edit_context = state.get("edit_context")
    if isinstance(edit_context, dict):
        manifest_path = Path(str(edit_context.get("manifest_path", ""))).resolve()
        template_pptx_path = Path(str(edit_context.get("template_pptx_path", ""))).resolve()
        if manifest_path.exists() and template_pptx_path.exists():
            return manifest_path, template_pptx_path
    return _resolve_template_package_paths(template_ref=None)


def _load_manifest(manifest_path: Path) -> dict[str, Any]:
    payload = _load_json_file(manifest_path)
    if int(payload.get("manifest_version", 0)) != MANIFEST_VERSION:
        raise ValueError(
            f"unsupported manifest version in {manifest_path}: {payload.get('manifest_version')}"
        )
    slides = payload.get("slides")
    if not isinstance(slides, list):
        raise ValueError(f"invalid manifest: missing slides list in {manifest_path}")
    return payload


def _build_fill_request(
    *,
    slide_index: int,
    field_specs: list[str] | None,
    manifest: dict[str, Any],
) -> dict[str, Any]:
    manifest_slide = _resolve_manifest_slide(
        manifest=manifest,
        slide_index=slide_index,
    )
    assignments = _validate_fill_fields(
        manifest_slide=manifest_slide,
        field_specs=field_specs,
    )
    return {
        "slide_name": str(manifest_slide["slide_name"]),
        "template_slide_index": int(manifest_slide["template_slide_index"]),
        "assignments": assignments,
    }


def _resolve_manifest_slide(
    *,
    manifest: dict[str, Any],
    slide_index: int,
) -> dict[str, Any]:
    manifest_slides = manifest["slides"]
    if slide_index < 0 or slide_index >= len(manifest_slides):
        raise ValueError(
            f"template slide index {slide_index} out of range; total template slides: {len(manifest_slides)}"
        )
    manifest_slide = manifest_slides[slide_index]
    if not isinstance(manifest_slide, dict):
        raise ValueError(f"invalid manifest slide entry at index {slide_index}")
    return manifest_slide


def _validate_fill_fields(
    *,
    manifest_slide: dict[str, Any],
    field_specs: list[str] | None,
) -> list[dict[str, Any]]:
    manifest_fields = manifest_slide.get("fields")
    if not isinstance(manifest_fields, list):
        raise ValueError(f"manifest slide has no fields: {manifest_slide.get('slide_name')}")

    if not manifest_fields:
        return []
    if not field_specs:
        raise ValueError("fill_template requires at least one --field; field input is required")

    fields_by_index: dict[int, dict[str, Any]] = {}
    for field in manifest_fields:
        if not isinstance(field, dict):
            continue
        index = field.get("index")
        if index is not None:
            fields_by_index[int(index)] = field

    assignments: list[dict[str, Any]] = []
    seen_indexes: set[int] = set()
    for raw_field_spec in field_specs:
        selected_index, raw_value = _parse_field_spec(raw_field_spec)
        if selected_index in seen_indexes:
            raise ValueError(
                f"duplicate field index on template slide `{manifest_slide['slide_name']}`: {selected_index}"
            )
        seen_indexes.add(selected_index)

        manifest_field = fields_by_index.get(selected_index)
        if manifest_field is None:
            raise ValueError(
                f"field index {selected_index} does not exist on template slide `{manifest_slide['slide_name']}`"
            )
        field_key = str(manifest_field.get("field_key") or manifest_field.get("description"))

        field_type = str(manifest_field.get("type", "")).strip()
        assignment: dict[str, Any] = {
            "field_key": field_key,
            "shape_id": int(manifest_field["shape_id"]),
            "type": field_type,
            "value": _validate_fill_value(
                slide_name=str(manifest_slide["slide_name"]),
                field_key=field_key,
                field_type=field_type,
                raw_value=raw_value,
            ),
        }
        assignments.append(assignment)

    missing_fields = [
        str(field.get("index"))
        for field in manifest_fields
        if int(field.get("index", -1)) not in seen_indexes
    ]
    if missing_fields:
        raise ValueError(
            f"missing fields on template slide `{manifest_slide['slide_name']}`: {', '.join(missing_fields)}"
        )
    return assignments


def _serialize_template_field(field: dict[str, Any]) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "index": int(field["index"]),
        "field_key": str(field.get("field_key") or ""),
        "description": str(field.get("description") or ""),
        "type": str(field.get("type") or ""),
        "shape_id": int(field["shape_id"]),
        "role_hint": field.get("role_hint"),
        "bbox": field.get("bbox"),
    }
    if field.get("text_excerpt") is not None:
        payload["text_excerpt"] = field["text_excerpt"]
    if field.get("image_name") is not None:
        payload["image_name"] = field["image_name"]
    return payload


def _validate_fill_value(
    *,
    slide_name: str,
    field_key: str,
    field_type: str,
    raw_value: Any,
) -> str:
    if not isinstance(raw_value, str):
        raise ValueError(
            f"field `{field_key}` on slide `{slide_name}` expects a string value for type `{field_type}`"
        )
    if field_type == "image":
        image_path = Path(raw_value).expanduser().resolve()
        if not image_path.exists():
            raise ValueError(
                f"image field `{field_key}` on slide `{slide_name}` points to a missing file: {image_path}"
            )
        return str(image_path)
    if field_type != "text":
        raise ValueError(
            f"unsupported field type `{field_type}` on slide `{slide_name}`; only text/image are supported"
        )
    return raw_value


def _remove_all_slides(presentation: Presentation) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _clone_slide(
    *,
    source_slide: Any,
    destination_presentation: Presentation,
) -> tuple[Any, dict[int, Any]]:
    destination_slide = destination_presentation.slides.add_slide(
        _select_destination_layout(
            source_slide=source_slide,
            destination_presentation=destination_presentation,
        )
    )
    _clear_slide_shapes(destination_slide)
    shape_map: dict[int, Any] = {}

    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        _rehydrate_shape_media_relationships(
            shape_element=new_element,
            source_slide=source_slide,
            destination_slide=destination_slide,
        )
        destination_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    for shape in destination_slide.shapes:
        shape_map.setdefault(int(shape.shape_id), shape)

    return destination_slide, shape_map


def _select_destination_layout(
    *,
    source_slide: Any,
    destination_presentation: Presentation,
) -> Any:
    source_layout = source_slide.slide_layout
    source_partname = str(source_layout.part.partname)
    for layout in destination_presentation.slide_layouts:
        if str(layout.part.partname) == source_partname:
            return layout

    source_layout_name = str(getattr(source_layout, "name", "")).strip()
    if source_layout_name:
        for layout in destination_presentation.slide_layouts:
            if str(getattr(layout, "name", "")).strip() == source_layout_name:
                return layout

    return _select_blank_layout(destination_presentation)


def _rehydrate_shape_media_relationships(
    *,
    shape_element: Any,
    source_slide: Any,
    destination_slide: Any,
) -> None:
    relationship_map: dict[str, str] = {}
    for blip in shape_element.iter(qn("a:blip")):
        previous_rel_id = blip.get(qn("r:embed"))
        if not previous_rel_id:
            continue
        next_rel_id = relationship_map.get(previous_rel_id)
        if next_rel_id is None:
            image_part = source_slide.part.related_part(previous_rel_id)
            if image_part is None:
                raise ValueError(f"unable to resolve image relationship `{previous_rel_id}` on source slide")
            _, next_rel_id = destination_slide.part.get_or_add_image_part(
                io.BytesIO(image_part.blob)
            )
            relationship_map[previous_rel_id] = next_rel_id
        blip.set(qn("r:embed"), next_rel_id)


def _select_blank_layout(presentation: Presentation) -> Any:
    layouts = presentation.slide_layouts
    if len(layouts) > 6:
        return layouts[6]
    return layouts[len(layouts) - 1]


def _clear_slide_shapes(slide: Any) -> None:
    for shape in list(slide.shapes):
        element = shape.element
        element.getparent().remove(element)


def _apply_slide_fields(
    *,
    slide: Any,
    shape_map: dict[int, Any],
    assignments: list[dict[str, Any]],
) -> None:
    for assignment in assignments:
        shape_id = int(assignment["shape_id"])
        shape = shape_map.get(shape_id)
        if shape is None:
            raise ValueError(f"shape_id {shape_id} not found when filling slide")
        if assignment["type"] == "text":
            _replace_text(shape=shape, value=str(assignment["value"]))
            continue
        if assignment["type"] == "image":
            replacement_shape = _replace_image(
                slide=slide,
                shape=shape,
                image_path=Path(str(assignment["value"])),
            )
            shape_map[shape_id] = replacement_shape
            continue
        raise ValueError(f"unsupported field type: {assignment['type']}")


def _replace_text(*, shape: Any, value: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"shape {shape.shape_id} does not support text replacement")
    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    if not paragraphs:
        raise ValueError(f"shape {shape.shape_id} has no paragraphs for text replacement")

    if "\n" not in value:
        _replace_text_in_first_run(text_frame=text_frame, value=value)
        return

    replacement_lines = value.split("\n")
    if not replacement_lines:
        replacement_lines = [""]
    _replace_text_in_paragraphs(text_frame=text_frame, replacement_lines=replacement_lines)


def _replace_text_in_first_run(*, text_frame: Any, value: str) -> None:
    first_paragraph = text_frame.paragraphs[0]
    _replace_paragraph_text_in_first_run(paragraph=first_paragraph, value=value)

    while len(text_frame.paragraphs) > 1:
        _remove_paragraph(text_frame.paragraphs[len(text_frame.paragraphs) - 1])


def _replace_text_in_paragraphs(*, text_frame: Any, replacement_lines: list[str]) -> None:
    paragraph_snapshots = [
        _capture_paragraph_snapshot(paragraph) for paragraph in text_frame.paragraphs
    ]
    existing_count = len(text_frame.paragraphs)
    shared_count = min(existing_count, len(replacement_lines))

    for index in range(shared_count):
        _replace_paragraph_text_in_first_run(
            paragraph=text_frame.paragraphs[index],
            value=replacement_lines[index],
        )

    while len(text_frame.paragraphs) > len(replacement_lines):
        _remove_paragraph(text_frame.paragraphs[len(text_frame.paragraphs) - 1])

    for index in range(shared_count, len(replacement_lines)):
        paragraph = text_frame.add_paragraph()
        snapshot = paragraph_snapshots[min(index, len(paragraph_snapshots) - 1)]
        _apply_paragraph_snapshot(
            paragraph=paragraph,
            snapshot=snapshot,
            text=replacement_lines[index],
        )


def _replace_paragraph_text_in_first_run(*, paragraph: Any, value: str) -> None:
    base_run = _resolve_base_run(paragraph)
    base_run_element = _get_run_element(base_run)
    base_run.text = value

    for run in list(paragraph.runs):
        if _get_run_element(run) is base_run_element:
            continue
        _remove_run(run)


def _resolve_base_run(paragraph: Any) -> Any:
    for run in paragraph.runs:
        if run.text:
            return run
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()


def _capture_paragraph_snapshot(paragraph: Any) -> dict[str, Any]:
    run_font_source = paragraph.runs[0].font if paragraph.runs else paragraph.font
    return {
        "alignment": paragraph.alignment,
        "level": paragraph.level,
        "font": _capture_font_snapshot(run_font_source),
    }


def _capture_font_snapshot(font: Any) -> dict[str, Any]:
    color = getattr(font, "color", None)
    return {
        "name": font.name,
        "size": font.size,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color_rgb": getattr(color, "rgb", None) if color is not None else None,
        "color_theme": getattr(color, "theme_color", None) if color is not None else None,
        "color_brightness": getattr(color, "brightness", None) if color is not None else None,
    }


def _apply_paragraph_snapshot(*, paragraph: Any, snapshot: dict[str, Any], text: str) -> None:
    paragraph.clear()
    paragraph.alignment = snapshot.get("alignment")
    paragraph.level = int(snapshot.get("level", 0) or 0)
    run = paragraph.add_run()
    run.text = text
    _apply_font_snapshot(font=run.font, snapshot=snapshot["font"])


def _apply_font_snapshot(*, font: Any, snapshot: dict[str, Any]) -> None:
    font.name = snapshot.get("name")
    font.size = snapshot.get("size")
    font.bold = snapshot.get("bold")
    font.italic = snapshot.get("italic")
    font.underline = snapshot.get("underline")

    color_rgb = snapshot.get("color_rgb")
    color_theme = snapshot.get("color_theme")
    color_brightness = snapshot.get("color_brightness")
    if color_rgb is not None:
        font.color.rgb = color_rgb
    elif color_theme is not None:
        font.color.theme_color = color_theme
        if color_brightness is not None:
            font.color.brightness = color_brightness


def _remove_paragraph(paragraph: Any) -> None:
    element = paragraph._element
    element.getparent().remove(element)


def _remove_run(run: Any) -> None:
    element = _get_run_element(run)
    if element is None:
        raise ValueError("run does not expose an XML element for removal")
    element.getparent().remove(element)


def _get_run_element(run: Any) -> Any:
    element = getattr(run, "_r", None)
    if element is not None:
        return element
    return getattr(run, "element", None)


def _replace_image(
    *,
    slide: Any,
    shape: Any,
    image_path: Path,
) -> Any:
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height
    shape.element.getparent().remove(shape.element)
    return slide.shapes.add_picture(str(image_path), left, top, width, height)


def _is_picture_shape(shape: Any) -> bool:
    return hasattr(shape, "image")


def _build_field_key(*, description: str, used_keys: set[str]) -> str:
    normalized = unicodedata.normalize("NFKD", description)
    ascii_only = normalized.encode("ascii", "ignore").decode("ascii").lower()
    slug = re.sub(r"[^a-z0-9]+", "_", ascii_only).strip("_")
    if not slug:
        slug = "field"
    candidate = slug
    suffix = 2
    while candidate in used_keys:
        candidate = f"{slug}_{suffix}"
        suffix += 1
    used_keys.add(candidate)
    return candidate


def _load_template_draft(draft_path: Path) -> dict[str, Any]:
    payload = _load_json_file(draft_path)
    if int(payload.get("draft_version", 0)) != TEMPLATE_DRAFT_VERSION:
        raise ValueError(
            f"unsupported template draft version in {draft_path}: {payload.get('draft_version')}"
        )
    return payload


def _parse_field_spec(value: str) -> tuple[int, str]:
    raw_value = value.strip()
    if ":" not in raw_value:
        raise ValueError(f'invalid --field value `{value}`; expected "index:value"')
    raw_index, raw_description = raw_value.split(":", 1)
    try:
        index = int(raw_index.strip())
    except ValueError as exc:
        raise ValueError(f'invalid --field index in `{value}`') from exc
    if index <= 0:
        raise ValueError(f"field index must be >= 1: {index}")
    description = raw_description.strip()
    if not description:
        raise ValueError(f"field value cannot be empty: `{value}`")
    return index, description


def _load_json_file(path: Path) -> dict[str, Any]:
    resolved = path.resolve()
    if not resolved.exists():
        raise ValueError(f"file does not exist: {resolved}")
    payload = json.loads(resolved.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError(f"expected JSON object in file: {resolved}")
    return payload


def _load_session_metadata(*, required: bool = False) -> dict[str, Any]:
    state_file = resolve_state_file_path()
    state = load_session_state_if_exists(state_file)
    if state is None:
        if required:
            raise ValueError(
                f"no active session state file found: {state_file}. Start with `pptxcli template create --from <file> --name <name>` or `pptxcli edit create --template <name> --output <file>` first."
            )
        return {}
    return state


def _persist_workflow_state(
    *,
    mode: str,
    active_template: str | None = None,
    edit_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return update_session_mode(
        resolve_state_file_path(),
        mode=mode,
        active_template=active_template,
        edit_context=edit_context,
    )


def _require_edit_context(state: dict[str, Any]) -> dict[str, Any]:
    mode = str(state.get("mode", "idle")).strip() or "idle"
    edit_context = state.get("edit_context")
    if mode != "edit_ppt" or not isinstance(edit_context, dict):
        raise ValueError(
            "no active edit draft found. Run `pptxcli edit create --template <name> --output <file>` first."
        )
    return edit_context


def _derive_template_name(manifest_path: Path) -> str:
    return normalize_template_name(manifest_path.parent.name)


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def _now_iso() -> str:
    return datetime.now(UTC).isoformat()
