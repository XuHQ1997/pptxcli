from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .models import BBox, Candidate, SlideCandidates, SlideObject, SlideObjects


def load_presentation(input_path: Path) -> object:
    return Presentation(str(input_path))


def inspect_slide(
    input_path: Path,
    slide_index: int,
    *,
    presentation: object | None = None,
) -> SlideCandidates:
    active_presentation = presentation or load_presentation(input_path)
    return inspect_slide_from_presentation(
        presentation=active_presentation,
        slide_index=slide_index,
    )


def inspect_slide_from_presentation(
    *,
    presentation: object,
    slide_index: int,
) -> SlideCandidates:
    slide = _get_slide(presentation=presentation, slide_index=slide_index)
    candidates: list[Candidate] = []

    for shape in slide.shapes:
        candidate = _shape_to_candidate(shape=shape, next_index=len(candidates) + 1)
        if candidate is not None:
            candidates.append(candidate)

    return SlideCandidates(
        slide_index=slide_index,
        slide_width=presentation.slide_width,
        slide_height=presentation.slide_height,
        candidates=candidates,
    )


def inspect_slide_objects(
    input_path: Path,
    slide_index: int,
    *,
    presentation: object | None = None,
) -> SlideObjects:
    active_presentation = presentation or load_presentation(input_path)
    return inspect_slide_objects_from_presentation(
        presentation=active_presentation,
        slide_index=slide_index,
    )


def inspect_slide_objects_from_presentation(
    *,
    presentation: object,
    slide_index: int,
) -> SlideObjects:
    slide = _get_slide(presentation=presentation, slide_index=slide_index)
    objects: list[SlideObject] = []

    for shape in slide.shapes:
        objects.append(_shape_to_object(shape=shape, next_index=len(objects) + 1))

    return SlideObjects(
        slide_index=slide_index,
        slide_width=presentation.slide_width,
        slide_height=presentation.slide_height,
        objects=objects,
    )


def _get_slide(*, presentation: object, slide_index: int) -> object:
    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise ValueError(
            f"slide index {slide_index} out of range; total slides: {len(presentation.slides)}"
        )
    return presentation.slides[slide_index]


def build_inspect_payload(
    *,
    command_name: str,
    input_path: Path,
    slide_index: int,
    output_path: Path | None = None,
    presentation: object | None = None,
) -> dict[str, object]:
    slide_data = inspect_slide_objects(
        input_path=input_path,
        slide_index=slide_index,
        presentation=presentation,
    )
    payload = slide_data.to_dict()
    payload.update(
        {
            "command": command_name,
            "input": str(input_path),
        }
    )

    if output_path is not None:
        resolved_output = output_path.resolve()
        resolved_output.parent.mkdir(parents=True, exist_ok=True)
        resolved_output.write_text(
            json.dumps(payload, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
        payload["output_path"] = str(resolved_output)

    return payload


def build_candidates_payload(
    *,
    command_name: str,
    input_path: Path,
    slide_index: int,
    presentation: object | None = None,
) -> dict[str, object]:
    slide_data = inspect_slide(
        input_path=input_path,
        slide_index=slide_index,
        presentation=presentation,
    )
    payload = slide_data.to_dict()
    payload.update(
        {
            "command": command_name,
            "input": str(input_path),
        }
    )
    return payload


def cmd_inspect(
    *,
    input_path: Path,
    slide_index: int,
    output_path: Path | None,
    command_name: str = "inspect",
    presentation: object | None = None,
) -> int:
    payload = build_inspect_payload(
        command_name=command_name,
        input_path=input_path,
        slide_index=slide_index,
        output_path=output_path,
        presentation=presentation,
    )
    print(json.dumps(payload, indent=2, ensure_ascii=False))
    return 0


def _shape_to_candidate(shape: object, next_index: int) -> Candidate | None:
    bbox = BBox(x=shape.left, y=shape.top, w=shape.width, h=shape.height)
    shape_id = int(shape.shape_id)
    candidate_id = f"slide-shape-{shape_id}"

    if _is_picture(shape):
        image = getattr(shape, "image", None)
        image_name = getattr(image, "filename", None)
        return Candidate(
            index=next_index,
            candidate_id=candidate_id,
            shape_id=shape_id,
            kind="image",
            role_hint="image",
            bbox=bbox,
            image_name=image_name,
        )

    if getattr(shape, "has_text_frame", False):
        text = _normalize_text(shape.text)
        if text:
            return Candidate(
                index=next_index,
                candidate_id=candidate_id,
                shape_id=shape_id,
                kind="text",
                role_hint=_infer_text_role_hint(text),
                bbox=bbox,
                text_excerpt=text[:80],
            )

    return None


def _shape_to_object(shape: object, next_index: int) -> SlideObject:
    bbox = BBox(x=shape.left, y=shape.top, w=shape.width, h=shape.height)
    image = getattr(shape, "image", None)
    image_name = getattr(image, "filename", None)
    text = _normalize_text(getattr(shape, "text", "")) if getattr(shape, "has_text_frame", False) else None
    object_type = _infer_object_type(shape=shape, text=text)
    placeholder_type = None
    if getattr(shape, "is_placeholder", False):
        placeholder_format = getattr(shape, "placeholder_format", None)
        placeholder_type = getattr(getattr(placeholder_format, "type", None), "name", None)

    return SlideObject(
        index=next_index,
        shape_id=int(shape.shape_id),
        object_type=object_type,
        shape_type=getattr(getattr(shape, "shape_type", None), "name", str(getattr(shape, "shape_type", "unknown"))),
        bbox=bbox,
        name=getattr(shape, "name", None),
        text=text or None,
        image_name=image_name,
        is_placeholder=bool(getattr(shape, "is_placeholder", False)),
        placeholder_type=placeholder_type,
    )


def _infer_object_type(*, shape: object, text: str | None) -> str:
    if _is_picture(shape):
        return "image"
    if text:
        return "text"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_table", False):
        return "table"
    return "shape"


def _is_picture(shape: object) -> bool:
    shape_type = getattr(shape, "shape_type", None)
    return shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image")


def _normalize_text(text: str) -> str:
    return " ".join(text.split())


def _infer_text_role_hint(text: str) -> str:
    lower_text = text.lower()
    if len(text) <= 40:
        return "title"
    if any(token in lower_text for token in ("summary", "overview", "agenda")):
        return "section_heading"
    if len(text) >= 120:
        return "body"
    return "text"
